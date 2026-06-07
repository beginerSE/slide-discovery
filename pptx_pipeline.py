"""PPTX text extraction + LibreOffice PNG thumbnail generation."""
from __future__ import annotations

import logging
import re
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

log = logging.getLogger("ingest.pptx")


@dataclass
class SlideExtract:
    page_no: int
    title: str
    body_text: str
    shape_types: list[str] = field(default_factory=list)
    has_chart: bool = False
    has_table: bool = False
    has_picture: bool = False
    thumbnail_path: Path | None = None


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def extract_slides(pptx_path: Path) -> list[SlideExtract]:
    prs = Presentation(str(pptx_path))
    out: list[SlideExtract] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        bodies: list[str] = []
        shape_types: list[str] = []
        has_chart = has_table = has_picture = False
        for shape in slide.shapes:
            try:
                stype = shape.shape_type
            except Exception:
                stype = None
            if stype is not None:
                shape_types.append(str(stype).split(".")[-1])
            if shape.has_chart:
                has_chart = True
            if stype == MSO_SHAPE_TYPE.TABLE or getattr(shape, "has_table", False):
                has_table = True
                try:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            t = cell.text.strip()
                            if t:
                                bodies.append(t)
                except Exception:
                    pass
            if stype == MSO_SHAPE_TYPE.PICTURE:
                has_picture = True
            txt = _shape_text(shape)
            if not txt:
                continue
            if not title and shape == slide.shapes.title:
                title = txt
            else:
                bodies.append(txt)
        if not title and bodies:
            first = bodies[0].splitlines()[0]
            if len(first) < 80:
                title = first
        body_text = "\n".join(bodies).strip()
        body_text = re.sub(r"\n{3,}", "\n\n", body_text)[:4000]
        out.append(
            SlideExtract(
                page_no=idx,
                title=title or "(無題)",
                body_text=body_text,
                shape_types=sorted(set(shape_types)),
                has_chart=has_chart,
                has_table=has_table,
                has_picture=has_picture,
            )
        )
    return out


def render_thumbnails(pptx_path: Path, out_dir: Path, dpi: int = 110) -> list[Path]:
    """Convert PPTX to per-page PNG using LibreOffice + pdftoppm.

    Returns sorted list of PNG paths (one per slide).
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    log.info("render start: PPTX->PDF %s (dpi=%d)", pptx_path.name, dpi)
    with tempfile.TemporaryDirectory(prefix="pptx_") as tmp:
        tmp_dir = Path(tmp)
        # 1) PPTX → PDF
        result = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(tmp_dir),
                str(pptx_path),
            ],
            capture_output=True,
            text=True,
            timeout=240,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"LibreOffice conversion failed: {result.stderr or result.stdout}"
            )
        pdfs = list(tmp_dir.glob("*.pdf"))
        if not pdfs:
            raise RuntimeError("LibreOffice produced no PDF output")
        pdf = pdfs[0]
        log.info("render: PDF->PNG %s", pptx_path.name)
        # 2) PDF → PNG (one per page). Write directly into out_dir (not the temp
        # dir) so a caller can observe pages landing on disk for live progress —
        # the PPTX→PDF (soffice) step above produces no intermediate files, so
        # the count staying at 0 during it cleanly signals "still converting".
        page_prefix = out_dir / "page"
        r2 = subprocess.run(
            ["pdftoppm", "-png", "-r", str(dpi), str(pdf), str(page_prefix)],
            capture_output=True,
            text=True,
            timeout=240,
        )
        if r2.returncode != 0:
            raise RuntimeError(
                f"pdftoppm failed: {r2.stderr or r2.stdout}"
            )
        pages = sorted(out_dir.glob("page-*.png"))
        if not pages:
            raise RuntimeError("pdftoppm produced no PNG output")
        # Rename in place to 1.png, 2.png, ... (same dir, so the total *.png
        # count stays stable while a progress poller is watching).
        final: list[Path] = []
        for i, p in enumerate(pages, start=1):
            target = out_dir / f"{i}.png"
            p.rename(target)
            final.append(target)
        log.info("render done: %s -> %d pages", pptx_path.name, len(final))
        return final
